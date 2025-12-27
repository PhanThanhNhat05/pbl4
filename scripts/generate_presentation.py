#!/usr/bin/env python3
"""
Assemble a PowerPoint presentation from images in slides_images/
Saves output to slides_images/presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

OUT_DIR = Path("slides_images")
PPTX_PATH = OUT_DIR / "presentation.pptx"


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    subtitle_tf = slide.placeholders[1]
    title_tf.text = title
    subtitle_tf.text = subtitle


def add_image_slide(prs, image_path, title=None, notes=None):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    left = top = Inches(0.5)
    pic = slide.shapes.add_picture(str(image_path), left, top, width=Inches(9))
    if title:
        # add title textbox
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(9), Inches(0.6))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main():
    if not OUT_DIR.exists():
        raise SystemExit("slides_images/ not found. Run image generation first.")
    prs = Presentation()
    # Title slide
    add_title_slide(prs,
                    "Phân tích dữ liệu huấn luyện & Giao diện hệ thống",
                    "Dự án PBL — (tên nhóm) — (dd/mm/yyyy)")

    # Slide: examples (if exists)
    examples = OUT_DIR / "slide_3_examples.png"
    if examples.exists():
        add_image_slide(prs, examples,
                        title="Ví dụ sóng ECG theo nhãn",
                        notes="Mỗi biểu đồ là sóng ECG mẫu đại diện cho từng nhãn. Nêu ngắn: Normal, AFib, PVC, ST change, Noise.")

    # Slide: distribution
    dist = OUT_DIR / "slide_4_distribution.png"
    if dist.exists():
        add_image_slide(prs, dist,
                        title="Phân bố 5 nhãn (imbalanced)",
                        notes="Ví dụ phân bố: Normal nhiều, Noise ít. Thảo luận: cần augmentation/weighted loss.")

    # Individual waveform slides
    labels = ["Normal","AFib","PVC","ST_change","Noise"]
    for lbl in labels:
        img = OUT_DIR / f"waveform_{lbl}.png"
        if img.exists():
            notes = {
                "Normal":"Mô tả: nhãn Normal — sóng bình thường.",
                "AFib":"Mô tả: nhãn AFib — nhịp không đều, cần chú ý.",
                "PVC":"Mô tả: PVC — các spike ngoại tâm thất.",
                "ST_change":"Mô tả: Thay đổi đoạn ST — biểu hiện tổn thương/ischemia.",
                "Noise":"Mô tả: Nhiễu / Artifact — loại bỏ/tái xử lý trước huấn luyện."
            }.get(lbl, "")
            add_image_slide(prs, img, title=f"ECG mẫu — {lbl}", notes=notes)

    prs.save(PPTX_PATH)
    print(f"Saved presentation to {PPTX_PATH}")


if __name__ == "__main__":
    main()





