#!/usr/bin/env python3
"""
Generate a more professional PowerPoint presentation using existing slide images.
Output: slides_images/professional_presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT_DIR = Path("slides_images")
OUT_PPTX = OUT_DIR / "professional_presentation.pptx"

PRIMARY_COLOR = RGBColor(28, 100, 162)  # blue
ACCENT_COLOR = RGBColor(88, 160, 120)   # green
TEXT_COLOR = RGBColor(33, 33, 33)


def set_text_run(run, text, size=18, bold=False, color=TEXT_COLOR):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_footer(slide, text):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    set_text_run(r, text, size=10, bold=False, color=RGBColor(120,120,120))


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    # colored left bar
    left_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(1.2), Inches(7.5))  # MSO_SHAPE_RECTANGLE=1
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = PRIMARY_COLOR
    left_bar.line.fill.background()
    # Title
    tx = slide.shapes.add_textbox(Inches(1.4), Inches(1.4), Inches(8), Inches(1.6))
    p = tx.text_frame.paragraphs[0]
    set_text_run(p.add_run(), title, size=28, bold=True)
    p2 = tx.text_frame.add_paragraph()
    set_text_run(p2.add_run(), subtitle, size=14, bold=False, color=RGBColor(80,80,80))
    add_footer(slide, "PBL Project — Generated slides")


def add_two_column_slide(prs, image_path, heading, bullets, notes=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    # left: image
    left = Inches(0.5)
    top = Inches(0.7)
    img_w = Inches(5.2)
    slide.shapes.add_picture(str(image_path), left, top, width=img_w)
    # right: text box
    tx = slide.shapes.add_textbox(Inches(6.0), Inches(0.8), Inches(3.5), Inches(5.5))
    tf = tx.text_frame
    set_text_run(tf.paragraphs[0].add_run(), heading, size=20, bold=True)
    for b in bullets:
        p = tf.add_paragraph()
        p.level = 1
        set_text_run(p.add_run(), f"• {b}", size=13, bold=False, color=RGBColor(60,60,60))
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_footer(slide, "PBL Project — Data & Interface")


def add_full_image_slide(prs, image_path, title=None, notes=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(0.7), width=Inches(9.0))
    if title:
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(9.0), Inches(0.6))
        set_text_run(tx.text_frame.paragraphs[0].add_run(), title, size=18, bold=True)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_footer(slide, "PBL Project — Visualization")


def main():
    if not OUT_DIR.exists():
        raise SystemExit("slides_images/ not found")
    prs = Presentation()
    # Title
    add_title_slide(prs,
                    "Phân tích dữ liệu huấn luyện & Giao diện hệ thống",
                    "Tổng quan dữ liệu, phân bố nhãn, và giao diện web")

    # Pipeline / Overview slide (use slide_3_examples if exists)
    examples = OUT_DIR / "slide_3_examples.png"
    if examples.exists():
        add_two_column_slide(prs, examples,
                             "Ví dụ sóng ECG theo nhãn",
                             ["Normal, AFib, PVC, ST change, Noise — mỗi nhãn 1 mẫu"],
                             notes="Giải thích ngắn từng nhãn, nêu đặc trưng.")

    # Distribution slide (large)
    dist = OUT_DIR / "class_distribution.png"
    if dist.exists():
        add_full_image_slide(prs, dist,
                             title="Phân bố 5 nhãn (ví dụ imbalanced)",
                             notes="Gợi ý: áp dụng augmentation, weighted loss, hoặc balanced sampling.")

    # Individual slides (image left, description right)
    label_notes = {
        "Normal":"Sóng ECG bình thường: nhịp đều, QRS rõ.",
        "AFib":"Rung nhĩ: nhịp không đều, absence of P waves.",
        "PVC":"Ngoại tâm thu: spike rời rạc, có thể khiến độ chính xác giảm.",
        "ST_change":"Thay đổi đoạn ST: cảnh báo tổn thương/ischemia.",
        "Noise":"Nhiễu/Artifact: cần lọc/loại bỏ trước huấn luyện."
    }
    for lbl in ["Normal","AFib","PVC","ST_change","Noise"]:
        img = OUT_DIR / f"waveform_{lbl}.png"
        if img.exists():
            add_two_column_slide(prs, img,
                                 heading=f"ECG mẫu — {lbl}",
                                 bullets=[label_notes.get(lbl, "")],
                                 notes=f"{lbl}: {label_notes.get(lbl, '')}")

    # UI overview slide (screenshot placeholders)
    add_two_column_slide(prs,
                         image_path=OUT_DIR / "slide_4_distribution.png" if (OUT_DIR / "slide_4_distribution.png").exists() else (OUT_DIR / "class_distribution.png"),
                         heading="Giao diện web — các trang chính",
                         bullets=["Login/Register", "Dashboard: tổng quan số liệu", "Measurement: thu thập / upload ECG", "History & Admin"],
                         notes="Trang chính theo repo: Login, Dashboard, Measurement, History, Admin.")

    # Recommendations / Next steps
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(9.0), Inches(1.2))
    set_text_run(tx.text_frame.paragraphs[0].add_run(), "Kết luận & Hướng tiếp theo", size=20, bold=True)
    bullets = [
        "Thu thêm mẫu cho các nhãn ít dữ liệu (Noise, ST_change).",
        "Áp dụng augmentation và class-weight khi huấn luyện.",
        "Triển khai inference real-time trên web (Measurement).",
        "Đánh giá model theo recall/precision trên từng nhãn."
    ]
    tx2 = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(8.6), Inches(4.0))
    tf = tx2.text_frame
    for b in bullets:
        p = tf.add_paragraph()
        set_text_run(p.add_run(), f"• {b}", size=14, bold=False, color=RGBColor(50,50,50))
    add_footer(slide, "PBL Project — Recommendations")

    prs.save(OUT_PPTX)
    print(f"Saved professional presentation to {OUT_PPTX}")


if __name__ == "__main__":
    main()





