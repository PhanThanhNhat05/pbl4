#!/usr/bin/env python3
"""
Insert a detailed 'Tổng quan dữ liệu' slide into slides_images/professional_presentation.pptx.
Overwrites the existing presentation file.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPTX_PATH = Path("slides_images/professional_presentation.pptx")
DIST_IMG = Path("slides_images/class_distribution.png")

def set_run(run, text, size=14, bold=False, color=RGBColor(33,33,33)):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_data_overview_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(0.8))
    t = title_box.text_frame.paragraphs[0]
    set_run(t.add_run(), "Tổng quan dữ liệu (chi tiết)", size=20, bold=True)

    # Left: bullets with specifics
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.2), Inches(4.8))
    tf = left_box.text_frame
    bullets = [
        "Tổng số mẫu (ví dụ): 1000 mẫu",
        "Sampling rate: 250 Hz (giả định trong ví dụ)",
        "Độ dài mỗi mẫu: 12 giây → ~3000 mẫu/chuỗi (12s × 250Hz)",
        "Kênh: 1 (single-lead / single-channel)",
        "Tiền xử lý: band-pass filter, loại bỏ baseline, normalization, segment/trim, resampling",
        "Gán nhãn: manual / từ metadata (5 nhãn: Normal, AFib, PVC, ST_change, Noise)",
        "Train/Val/Test (ví dụ): 70% / 15% / 15% → 700 / 150 / 150",
        "Vấn đề chính: class imbalance → ảnh hưởng tới recall trên nhãn ít mẫu"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.level = 0
        set_run(p.add_run(), f"• {b}", size=13, bold=False, color=RGBColor(50,50,50))

    # Right: class distribution image if exists
    if DIST_IMG.exists():
        slide.shapes.add_picture(str(DIST_IMG), Inches(6.2), Inches(1.6), width=Inches(3.4))
    else:
        # small note if image missing
        note_box = slide.shapes.add_textbox(Inches(6.2), Inches(1.6), Inches(3.4), Inches(0.6))
        set_run(note_box.text_frame.paragraphs[0].add_run(), "Class distribution image not found", size=12, bold=False)

    # Speaker notes (Vietnamese)
    notes = (
        "Nói: Slide này cung cấp con số và thông số chính của dataset mẫu.\n"
        "- Nhấn mạnh sampling rate và độ dài mẫu (12s) vì ảnh hưởng tới preprocessing và model input shape.\n"
        "- Trình bày split train/val/test và nêu ví dụ counts; giải thích rủi ro do imbalance và bước xử lý dự kiến."
    )
    slide.notes_slide.notes_text_frame.text = notes


def main():
    if not PPTX_PATH.exists():
        raise SystemExit(f"{PPTX_PATH} not found — hãy chạy generate_presentation_pro.py trước.")
    prs = Presentation(str(PPTX_PATH))
    add_data_overview_slide(prs)
    prs.save(str(PPTX_PATH))
    print(f"Updated {PPTX_PATH} with data overview slide.")


if __name__ == "__main__":
    main()





